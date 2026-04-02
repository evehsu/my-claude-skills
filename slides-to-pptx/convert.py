#!/usr/bin/env python3
"""Convert HTML slide decks to PPTX by screenshotting each slide."""

import argparse
import subprocess
import sys
import tempfile
from pathlib import Path


def ensure_playwright():
    try:
        import playwright  # noqa: F401
    except ImportError:
        print("Installing playwright...", flush=True)
        subprocess.check_call([sys.executable, "-m", "pip", "install", "playwright", "-q"])

    # Check that the browser binary expected by the installed playwright version actually exists
    from playwright.sync_api import sync_playwright
    with sync_playwright() as p:
        exe = p.chromium.executable_path
    if not Path(exe).exists():
        print("Downloading Chromium browser for playwright...", flush=True)
        subprocess.check_call([sys.executable, "-m", "playwright", "install", "chromium"])


def parse_notes(html_path: Path) -> list[str]:
    from bs4 import BeautifulSoup

    html = html_path.read_text(encoding="utf-8")
    soup = BeautifulSoup(html, "html.parser")
    notes = []
    for slide in soup.select("section.slide"):
        aside = slide.find("aside", class_="notes")
        notes.append(aside.get_text(strip=True) if aside else "")
    return notes


def screenshot_slides(html_path: Path, tmp_dir: Path) -> list[Path]:
    from playwright.sync_api import sync_playwright

    screenshots = []
    file_url = html_path.resolve().as_uri()

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        context = browser.new_context(
            viewport={"width": 1280, "height": 720},
            device_scale_factor=2,
        )
        page = context.new_page()
        page.goto(file_url, wait_until="networkidle")
        page.wait_for_timeout(500)

        slide_count = page.evaluate("document.querySelectorAll('.slide').length")
        print(f"Found {slide_count} slides", flush=True)

        # Navigate to first slide and screenshot
        page.keyboard.press("Home")
        page.wait_for_timeout(150)
        path = tmp_dir / "slide_000.png"
        page.screenshot(path=str(path))
        screenshots.append(path)

        for i in range(1, slide_count):
            page.keyboard.press("ArrowRight")
            page.wait_for_timeout(150)
            path = tmp_dir / f"slide_{i:03d}.png"
            page.screenshot(path=str(path))
            screenshots.append(path)

        browser.close()

    return sorted(screenshots)


def build_pptx(screenshots: list[Path], notes: list[str], output_path: Path):
    from pptx import Presentation
    from pptx.util import Emu

    # 10" x 5.625" in EMU (914400 EMU per inch)
    SLIDE_WIDTH = Emu(9144000)
    SLIDE_HEIGHT = Emu(5143500)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = prs.slide_layouts[6]  # Blank layout

    for i, screenshot_path in enumerate(screenshots):
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(screenshot_path),
            left=0,
            top=0,
            width=SLIDE_WIDTH,
            height=SLIDE_HEIGHT,
        )
        note_text = notes[i] if i < len(notes) else ""
        if note_text:
            slide.notes_slide.notes_text_frame.text = note_text

    prs.save(str(output_path))


def main():
    parser = argparse.ArgumentParser(description="Convert HTML slides to PPTX")
    parser.add_argument("input_html", help="Path to HTML slide deck")
    parser.add_argument("-o", "--output", help="Output PPTX path (default: same name as input)")
    args = parser.parse_args()

    html_path = Path(args.input_html).expanduser().resolve()
    if not html_path.exists():
        print(f"Error: File not found: {html_path}", file=sys.stderr)
        sys.exit(1)

    output_path = Path(args.output).expanduser().resolve() if args.output else html_path.with_suffix(".pptx")

    ensure_playwright()

    with tempfile.TemporaryDirectory() as tmp:
        tmp_dir = Path(tmp)

        print("Extracting speaker notes...", flush=True)
        notes = parse_notes(html_path)

        print("Taking screenshots...", flush=True)
        screenshots = screenshot_slides(html_path, tmp_dir)

        print(f"Building PPTX with {len(screenshots)} slides...", flush=True)
        build_pptx(screenshots, notes, output_path)

    print(f"Done: {output_path}")
    print(f"Open with: open \"{output_path}\"")


if __name__ == "__main__":
    main()
